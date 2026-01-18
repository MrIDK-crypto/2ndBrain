"""
Video Generator for Training Materials
Creates training videos from PowerPoint presentations
"""

import json
from pathlib import Path
from typing import Dict, List, Optional
from pptx import Presentation
from gtts import gTTS
from moviepy.editor import (
    ImageClip, AudioFileClip, concatenate_videoclips,
    CompositeVideoClip, TextClip
)
from PIL import Image, ImageDraw, ImageFont
import io
import os


class VideoGenerator:
    """Generate training videos from presentations"""

    def __init__(self, output_dir: str):
        """
        Initialize video generator

        Args:
            output_dir: Directory to save generated videos
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

        self.temp_dir = self.output_dir / "temp"
        self.temp_dir.mkdir(parents=True, exist_ok=True)

    def extract_ppt_content(self, ppt_path: str) -> List[Dict]:
        """
        Extract content from PowerPoint presentation

        Args:
            ppt_path: Path to PowerPoint file

        Returns:
            List of slide dictionaries
        """
        prs = Presentation(ppt_path)
        slides_content = []

        for i, slide in enumerate(prs.slides):
            slide_data = {
                'slide_number': i + 1,
                'title': '',
                'content': [],
                'notes': ''
            }

            # Extract title
            if slide.shapes.title:
                slide_data['title'] = slide.shapes.title.text

            # Extract content
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    if shape != slide.shapes.title:
                        slide_data['content'].append(shape.text)

            # Extract speaker notes
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame:
                    slide_data['notes'] = notes_frame.text

            slides_content.append(slide_data)

        return slides_content

    def create_slide_image(
        self,
        slide_data: Dict,
        width: int = 1920,
        height: int = 1080
    ) -> str:
        """
        Create an image for a slide

        Args:
            slide_data: Slide content dictionary
            width: Image width
            height: Image height

        Returns:
            Path to created image
        """
        # Create image
        img = Image.new('RGB', (width, height), color='white')
        draw = ImageDraw.Draw(img)

        # Try to load fonts, fall back to default if not available
        try:
            title_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 72)
            content_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 48)
        except:
            title_font = ImageFont.load_default()
            content_font = ImageFont.load_default()

        # Draw title
        title = slide_data.get('title', '')
        if title:
            # Title background
            draw.rectangle([(0, 0), (width, 200)], fill=(0, 51, 102))
            # Title text
            draw.text((100, 80), title, fill='white', font=title_font)

        # Draw content bullets
        y_position = 300
        for content_text in slide_data.get('content', []):
            # Split long text into lines
            lines = self._wrap_text(content_text, content_font, width - 200)
            for line in lines:
                if y_position < height - 100:
                    draw.text((100, y_position), f"• {line}", fill=(64, 64, 64), font=content_font)
                    y_position += 80

        # Save image
        img_path = self.temp_dir / f"slide_{slide_data['slide_number']}.png"
        img.save(str(img_path))

        return str(img_path)

    def _wrap_text(self, text: str, font, max_width: int) -> List[str]:
        """Wrap text to fit within max width"""
        words = text.split()
        lines = []
        current_line = []

        for word in words:
            test_line = ' '.join(current_line + [word])
            # Simple character-based wrapping (proper would use font.getsize)
            if len(test_line) * 20 < max_width:  # Rough estimate
                current_line.append(word)
            else:
                if current_line:
                    lines.append(' '.join(current_line))
                current_line = [word]

        if current_line:
            lines.append(' '.join(current_line))

        return lines if lines else [text]

    def generate_narration(self, slide_data: Dict) -> str:
        """
        Generate audio narration for a slide

        Args:
            slide_data: Slide content

        Returns:
            Path to audio file
        """
        # Create narration text
        narration_parts = []

        # Add title
        title = slide_data.get('title', '')
        if title:
            narration_parts.append(title)

        # Add speaker notes or content
        notes = slide_data.get('notes', '')
        if notes:
            narration_parts.append(notes)
        else:
            # Fall back to content
            content = slide_data.get('content', [])
            for text in content:
                # Clean up bullet points
                clean_text = text.replace('•', '').strip()
                if clean_text:
                    narration_parts.append(clean_text)

        narration_text = '. '.join(narration_parts)

        if not narration_text.strip():
            narration_text = "This slide shows important information."

        # Generate speech
        audio_path = self.temp_dir / f"audio_{slide_data['slide_number']}.mp3"

        try:
            tts = gTTS(text=narration_text, lang='en', slow=False)
            tts.save(str(audio_path))
        except Exception as e:
            print(f"Error generating audio: {e}")
            # Create silent audio as fallback
            return None

        return str(audio_path)

    def create_video_from_ppt(
        self,
        ppt_path: str,
        output_path: str,
        fps: int = 24
    ):
        """
        Create video from PowerPoint presentation

        Args:
            ppt_path: Path to PowerPoint file
            output_path: Path to save video
            fps: Frames per second
        """
        print(f"Creating video from {Path(ppt_path).name}...")

        # Extract slides
        slides_content = self.extract_ppt_content(ppt_path)

        if not slides_content:
            print("  No slides found in presentation")
            return

        video_clips = []

        # Process each slide
        for slide_data in slides_content:
            print(f"  Processing slide {slide_data['slide_number']}...")

            # Create slide image
            img_path = self.create_slide_image(slide_data)

            # Generate narration
            audio_path = self.generate_narration(slide_data)

            # Create video clip
            if audio_path and os.path.exists(audio_path):
                # Get audio duration
                audio_clip = AudioFileClip(audio_path)
                duration = audio_clip.duration
                audio_clip.close()

                # Create image clip with audio duration
                img_clip = ImageClip(img_path, duration=duration)

                # Add audio
                audio = AudioFileClip(audio_path)
                video_clip = img_clip.set_audio(audio)
            else:
                # No audio, use fixed duration
                duration = 5  # 5 seconds per slide
                img_clip = ImageClip(img_path, duration=duration)
                video_clip = img_clip

            video_clips.append(video_clip)

        # Concatenate all clips
        if video_clips:
            print("  Combining video clips...")
            final_video = concatenate_videoclips(video_clips, method="compose")

            # Write video file
            output_file = Path(output_path)
            output_file.parent.mkdir(parents=True, exist_ok=True)

            print(f"  Writing video to {output_file}...")
            final_video.write_videofile(
                str(output_file),
                fps=fps,
                codec='libx264',
                audio_codec='aac',
                verbose=False,
                logger=None
            )

            # Clean up
            final_video.close()
            for clip in video_clips:
                clip.close()

            print(f"✓ Created video: {output_file}")
        else:
            print("  No video clips created")

    def generate_videos_from_presentations(self, presentations_dir: str):
        """
        Generate videos from all presentations in a directory

        Args:
            presentations_dir: Directory containing PowerPoint files
        """
        print("\nGenerating training videos from presentations...")

        ppt_dir = Path(presentations_dir)

        # Find all PowerPoint files
        ppt_files = list(ppt_dir.glob("*.pptx"))

        if not ppt_files:
            print("No PowerPoint files found")
            return

        print(f"Found {len(ppt_files)} presentations")

        # Generate video for each presentation
        for ppt_file in ppt_files[:3]:  # Limit to first 3 for testing
            video_name = ppt_file.stem + ".mp4"
            video_path = self.output_dir / video_name

            try:
                self.create_video_from_ppt(str(ppt_file), str(video_path))
            except Exception as e:
                print(f"  ✗ Error creating video for {ppt_file.name}: {e}")

        print("\n✓ Video generation complete")

    def cleanup_temp_files(self):
        """Remove temporary files"""
        import shutil
        if self.temp_dir.exists():
            shutil.rmtree(self.temp_dir)
            self.temp_dir.mkdir(parents=True, exist_ok=True)
        print("✓ Cleaned up temporary files")


if __name__ == "__main__":
    from config.config import Config

    generator = VideoGenerator(output_dir=str(Config.VIDEO_OUTPUT_DIR))

    # Generate videos from presentations
    generator.generate_videos_from_presentations(
        presentations_dir=str(Config.OUTPUT_DIR / "powerpoints")
    )

    # Cleanup
    generator.cleanup_temp_files()
